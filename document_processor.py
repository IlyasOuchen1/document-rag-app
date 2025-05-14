"""
document_processor.py: Handles document ingestion and text extraction
"""

import os
import tempfile
import PyPDF2
import pandas as pd
from docx import Document
from pptx import Presentation
import csv
import json
from typing import List, Dict, Any, Tuple
from langchain.text_splitter import RecursiveCharacterTextSplitter
import logging

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class DocumentProcessor:
    """Class for processing various document types and extracting text content."""
    
    def __init__(self, chunk_size: int = 500, chunk_overlap: int = 100):
        """
        Initialize the document processor.
        
        Args:
            chunk_size: The size of text chunks to create
            chunk_overlap: The overlap between consecutive chunks
        """
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap
        )
        self.supported_extensions = {
            '.pdf': self._process_pdf,
            '.docx': self._process_docx,
            '.pptx': self._process_pptx,
            '.csv': self._process_csv,
            '.txt': self._process_txt,
            '.json': self._process_json
        }
        
    def process_file(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Process a file and extract text content.
        
        Args:
            file_path: Path to the file to process
            
        Returns:
            List of dictionaries containing text chunks and metadata
        """
        ext = os.path.splitext(file_path)[1].lower()
        if ext not in self.supported_extensions:
            raise ValueError(f"Unsupported file type: {ext}")
        
        try:
            return self.supported_extensions[ext](file_path)
        except Exception as e:
            logger.error(f"Error processing file {file_path}: {str(e)}")
            raise
    
    def _process_pdf(self, file_path: str) -> List[Dict[str, Any]]:
        """Process PDF files and extract text."""
        chunks = []
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page_num, page in enumerate(pdf_reader.pages):
                text = page.extract_text()
                if text.strip():
                    chunks.append({
                        'text': text,
                        'metadata': {
                            'source': file_path,
                            'page': page_num + 1,
                            'type': 'pdf'
                        }
                    })
        return chunks
    
    def _process_docx(self, file_path: str) -> List[Dict[str, Any]]:
        """Process DOCX files and extract text."""
        chunks = []
        doc = Document(file_path)
        for para in doc.paragraphs:
            if para.text.strip():
                chunks.append({
                    'text': para.text,
                    'metadata': {
                        'source': file_path,
                        'type': 'docx'
                    }
                })
        return chunks
    
    def _process_pptx(self, file_path: str) -> List[Dict[str, Any]]:
        """Process PPTX files and extract text."""
        chunks = []
        prs = Presentation(file_path)
        for slide_num, slide in enumerate(prs.slides):
            text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
            if text.strip():
                chunks.append({
                    'text': text,
                    'metadata': {
                        'source': file_path,
                        'slide': slide_num + 1,
                        'type': 'pptx'
                    }
                })
        return chunks
    
    def _process_csv(self, file_path: str) -> List[Dict[str, Any]]:
        """Process CSV files and extract text."""
        chunks = []
        df = pd.read_csv(file_path)
        for _, row in df.iterrows():
            text = row.to_string()
            if text.strip():
                chunks.append({
                    'text': text,
                    'metadata': {
                        'source': file_path,
                        'type': 'csv'
                    }
                })
        return chunks
    
    def _process_txt(self, file_path: str) -> List[Dict[str, Any]]:
        """Process TXT files and extract text."""
        chunks = []
        with open(file_path, 'r', encoding='utf-8') as file:
            text = file.read()
            if text.strip():
                chunks.append({
                    'text': text,
                    'metadata': {
                        'source': file_path,
                        'type': 'txt'
                    }
                })
        return chunks
    
    def _process_json(self, file_path: str) -> List[Dict[str, Any]]:
        """Process JSON files and extract text."""
        chunks = []
        with open(file_path, 'r', encoding='utf-8') as file:
            data = json.load(file)
            text = json.dumps(data, indent=2)
            if text.strip():
                chunks.append({
                    'text': text,
                    'metadata': {
                        'source': file_path,
                        'type': 'json'
                    }
                })
        return chunks
    
    def _create_chunks(self, text: str, metadata: Dict[str, Any]) -> List[Dict[str, Any]]:
        """
        Split text into chunks and add metadata.
        
        Args:
            text: The text to split
            metadata: Metadata about the source document
            
        Returns:
            List of dictionaries with text chunks and metadata
        """
        chunks = self.text_splitter.split_text(text)
        
        return [
            {
                "text": chunk,
                "metadata": {
                    **metadata,
                    "chunk_id": i
                }
            }
            for i, chunk in enumerate(chunks)
        ]
    
    def process_uploaded_file(self, uploaded_file) -> List[Dict[str, Any]]:
        """
        Process an uploaded file object (e.g., from Streamlit).
        
        Args:
            uploaded_file: File object from upload widget
            
        Returns:
            List of dictionaries with text chunks and metadata
        """
        # Create a temporary file to process
        with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(uploaded_file.name)[1]) as temp_file:
            temp_file.write(uploaded_file.getvalue())
            temp_file_path = temp_file.name
        
        try:
            result = self.process_file(temp_file_path)
        finally:
            # Clean up the temporary file
            os.unlink(temp_file_path)
            
        return result

    def batch_process_directory(self, directory_path: str) -> List[Dict[str, Any]]:
        """
        Process all supported files in a directory.
        
        Args:
            directory_path: Path to directory containing documents
            
        Returns:
            List of dictionaries with text chunks and metadata
        """
        all_chunks = []
        for root, _, files in os.walk(directory_path):
            for file in files:
                file_path = os.path.join(root, file)
                try:
                    chunks = self.process_file(file_path)
                    all_chunks.extend(chunks)
                    print(f"Processed {file_path}: {len(chunks)} chunks created")
                except Exception as e:
                    logger.error(f"Error processing {file_path}: {str(e)}")
        
        return all_chunks